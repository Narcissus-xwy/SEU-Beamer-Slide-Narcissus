"""PDF to PPTX converter for SEU-Beamer-Slide-Narcissus.

Usage:
    python pdf2ppt.py <path-to-pdf> [--dpi 300]

The script automatically searches for the PDF in:
1. The given path as-is
2. The out/ directory next to this script (project root)

Or called from pdf2ppt.bat (double-click).
"""
import sys
import os
import argparse
from pathlib import Path

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import io


PDF_BASE_DPI = 72


def resolve_pdf_path(pdf_arg: str) -> Path:
    """Resolve PDF path, searching in project out/ if not found directly."""
    # Strip quotes
    pdf_arg = pdf_arg.strip('"\'')

    if not pdf_arg:
        raise FileNotFoundError("输入的路径为空")

    # Try as-is
    path = Path(pdf_arg)
    if path.suffix.lower() != ".pdf":
        path = path.with_suffix(".pdf")
    if path.exists():
        return path

    # Try in out/ directory (project root)
    script_dir = Path(__file__).resolve().parent
    out_path = script_dir / "out" / path.name
    if out_path.exists():
        print(f"在 out/ 目录找到文件: {out_path}")
        return out_path

    raise FileNotFoundError(
        f"文件不存在（已尝试: {path} 及 out/{path.name}）"
    )


def pdf_to_pptx(pdf_path: str, dpi: int = 1200) -> str:
    pdf_path = resolve_pdf_path(pdf_path)

    out_path = pdf_path.with_suffix(".pptx")
    # 确保输出文件名合理，不会产生意外文件
    if not out_path.name or out_path.name.startswith("."):
        raise ValueError(f"输出路径不合理: {out_path}")
    zoom = dpi / PDF_BASE_DPI

    doc = fitz.open(str(pdf_path))
    total = doc.page_count
    if total == 0:
        doc.close()
        raise ValueError("PDF 文件为空（0 页）")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    for i in range(total):
        page = doc[i]
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat)
        img_bytes = pix.tobytes("png")

        slide = prs.slides.add_slide(blank)
        for ph in list(slide.placeholders):
            ph.element.getparent().remove(ph.element)

        slide.shapes.add_picture(io.BytesIO(img_bytes), 0, 0, prs.slide_width, prs.slide_height)

    doc.close()
    prs.save(str(out_path))
    return str(out_path)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="PDF 转 PPTX 工具")
    parser.add_argument("pdf", help="PDF 文件路径")
    parser.add_argument("--dpi", type=int, default=1200,
                        help="渲染分辨率（默认 1200 DPI）")
    args = parser.parse_args()

    try:
        result = pdf_to_pptx(args.pdf, dpi=args.dpi)
        size_kb = os.path.getsize(result) / 1024
        print(f"\n转换完成！")
        print(f"DPI:     {args.dpi}")
        print(f"输出:    {result}")
        print(f"大小:    {round(size_kb / 1024, 1)} MB")
    except Exception as e:
        print(f"\n[错误] {e}")
        sys.exit(1)
